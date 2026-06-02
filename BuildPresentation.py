"""
Build the ThroughLine pitch deck following throughline_deck_prompt.md.

Visual identity:
  - Primary: dark navy (#0B1F3A) + teal accent (#19C2B5)
  - Headline (dark) slides: navy background, white text
  - Content slides: light ivory background (#F5F7FA), navy text, teal accents
  - Display font: "Inter" / fallback "Helvetica Neue" via python-pptx ("Calibri" if unavailable on system).
    We set font names explicitly; PPT will substitute gracefully.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ---------- Palette ----------
NAVY = RGBColor(0x0B, 0x1F, 0x3A)
NAVY_DEEP = RGBColor(0x07, 0x16, 0x2B)
TEAL = RGBColor(0x19, 0xC2, 0xB5)
TEAL_SOFT = RGBColor(0x7F, 0xE0, 0xD8)
IVORY = RGBColor(0xF5, 0xF7, 0xFA)
SLATE = RGBColor(0x2A, 0x3A, 0x55)
MUTED = RGBColor(0x5C, 0x6B, 0x82)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RULE = RGBColor(0xD8, 0xDF, 0xEA)
AMBER = RGBColor(0xF2, 0xA8, 0x3B)

# Fonts
HEAD_FONT = "Inter"
BODY_FONT = "Inter"

# Widescreen 16:9
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def add_rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp


def add_text(
    slide, x, y, w, h, text, *,
    font=BODY_FONT, size=14, bold=False, color=NAVY,
    align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.15,
):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb


def add_bullets(
    slide, x, y, w, h, items, *,
    font=BODY_FONT, size=16, color=NAVY, lead_color=TEAL,
    bullet_char="\u25A0", line_spacing=1.25, gap_pt=8,
):
    """
    items: list of dicts: {"lead": str (bold colored), "rest": str}
           or plain strings.
    """
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)

    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(gap_pt)

        # bullet glyph
        rb = p.add_run()
        rb.text = bullet_char + "  "
        rb.font.name = font
        rb.font.size = Pt(size)
        rb.font.bold = True
        rb.font.color.rgb = lead_color

        if isinstance(it, dict):
            lead = it.get("lead", "")
            rest = it.get("rest", "")
            if lead:
                rl = p.add_run()
                rl.text = lead
                rl.font.name = font
                rl.font.size = Pt(size)
                rl.font.bold = True
                rl.font.color.rgb = color
            if rest:
                rr = p.add_run()
                # If lead ends with punctuation, keep tight; else add separator
                sep = " " if not lead else ("" if lead.endswith((":", "—", "-")) else " — ")
                rr.text = sep + rest
                rr.font.name = font
                rr.font.size = Pt(size)
                rr.font.color.rgb = color
        else:
            r = p.add_run()
            r.text = it
            r.font.name = font
            r.font.size = Pt(size)
            r.font.color.rgb = color
    return tb


def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_footer(slide, idx, total, on_dark=False):
    txt_color = WHITE if on_dark else MUTED
    rule_color = TEAL if on_dark else RULE
    # left brand
    add_text(
        slide, Inches(0.5), Inches(7.05), Inches(4), Inches(0.3),
        "ThroughLine  ·  Dreamforce.io Pitch",
        font=BODY_FONT, size=10, color=txt_color, bold=True,
    )
    # right page no
    add_text(
        slide, Inches(11.5), Inches(7.05), Inches(1.4), Inches(0.3),
        f"{idx:02d} / {total:02d}",
        font=BODY_FONT, size=10, color=txt_color, align=PP_ALIGN.RIGHT,
    )
    # thin top rule (under header)
    add_rect(slide, Inches(0.5), Inches(7.0), Inches(12.333), Emu(9000), rule_color)


def add_header_bar(slide, kicker, *, on_dark=False):
    txt_color = TEAL if on_dark else TEAL
    # small teal square accent
    add_rect(slide, Inches(0.5), Inches(0.55), Inches(0.18), Inches(0.18), TEAL)
    add_text(
        slide, Inches(0.78), Inches(0.45), Inches(10), Inches(0.4),
        kicker.upper(),
        font=BODY_FONT, size=11, bold=True, color=txt_color,
    )


def add_headline(slide, text, *, on_dark=False, y=Inches(1.0), size=34):
    color = WHITE if on_dark else NAVY
    add_text(
        slide, Inches(0.5), y, Inches(12.333), Inches(1.6),
        text, font=HEAD_FONT, size=size, bold=True, color=color, line_spacing=1.1,
    )


# ============================================================
# Build deck
# ============================================================

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]

TOTAL = 16  # placeholder; final set after all slides added

# ------------------- SLIDE 1: Title -------------------
s = prs.slides.add_slide(blank)
set_bg(s, NAVY_DEEP)

# soft teal diagonal accent
acc = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.35), SLIDE_H)
acc.fill.solid(); acc.fill.fore_color.rgb = TEAL; acc.line.fill.background()

# tiny eyebrow
add_text(s, Inches(0.95), Inches(2.4), Inches(8), Inches(0.4),
         "DREAMFORCE.IO  ·  PRODUCT INTERN TAKE-HOME",
         font=BODY_FONT, size=12, bold=True, color=TEAL)

# Wordmark
add_text(s, Inches(0.9), Inches(2.85), Inches(12), Inches(1.6),
         "ThroughLine",
         font=HEAD_FONT, size=84, bold=True, color=WHITE, line_spacing=1.0)

# Underline accent
add_rect(s, Inches(0.95), Inches(4.05), Inches(1.6), Emu(36000), TEAL)

# One-line vision
add_text(s, Inches(0.95), Inches(4.25), Inches(11.5), Inches(1.2),
         "A career advisor that actually knows you — and is built for the student, not the recruiter.",
         font=HEAD_FONT, size=22, color=WHITE, line_spacing=1.25)

# Author/footer
add_text(s, Inches(0.95), Inches(6.5), Inches(8), Inches(0.4),
         "15–20 min pitch  ·  Product & Design review",
         font=BODY_FONT, size=11, color=TEAL_SOFT, bold=True)

# ------------------- SLIDE 2: User problem (Part 1) -------------------
s = prs.slides.add_slide(blank)
set_bg(s, IVORY)
add_header_bar(s, "The user problem  ·  Part 1 of 2")
add_headline(s, "Every career platform asks students to perform a finished version of themselves.")

add_bullets(
    s, Inches(0.5), Inches(2.5), Inches(7.6), Inches(4.2),
    [
        {"lead": "Built for the employer side.",
         "rest": "LinkedIn and Handshake connect users to jobs — they don't build the user into someone qualified for one."},
        {"lead": "Profiles assume experience the user doesn't have yet.",
         "rest": "A student with no work history gets little value out of starting a profile."},
        {"lead": "The advising is miscalibrated.",
         "rest": "It reads your resume, not your actual goals — so the recommendations drift toward generic."},
    ],
    size=17,
)

# right callout
add_rect(s, Inches(8.45), Inches(2.5), Inches(4.4), Inches(4.0), WHITE, RULE)
add_rect(s, Inches(8.45), Inches(2.5), Inches(0.12), Inches(4.0), TEAL)
add_text(s, Inches(8.7), Inches(2.7), Inches(4.1), Inches(0.5),
         "WHO IS UNDERSERVED", font=BODY_FONT, size=11, bold=True, color=TEAL)
add_text(s, Inches(8.7), Inches(3.1), Inches(4.1), Inches(3.2),
         "College students early in the search — including the ones who don't yet know what they want.\n\n"
         "They're the exact users existing platforms quietly assume away.",
         font=HEAD_FONT, size=16, color=NAVY, line_spacing=1.3)

# ------------------- SLIDE 3: User problem (Part 2) -------------------
s = prs.slides.add_slide(blank)
set_bg(s, IVORY)
add_header_bar(s, "The user problem  ·  Part 2 of 2")
add_headline(s, "The truest things about your career are the things you'd never put on a resume.")

add_bullets(
    s, Inches(0.5), Inches(2.5), Inches(12.3), Inches(4.2),
    [
        {"lead": "What you're building but haven't finished.",
         "rest": "The half-done project that signals more than the title at the top of a LinkedIn page."},
        {"lead": "What you're unsure about.",
         "rest": "Doubts about direction that any honest advisor would need to know first."},
        {"lead": "What other industries you're quietly considering.",
         "rest": "The optionality students hide from outward-facing profiles."},
        {"lead": "The cost of hiding it:",
         "rest": "every tool meets a version of you that isn't real, and gives advice tuned to that version."},
    ],
    size=17,
)

# ------------------- SLIDE 4: LinkedIn — strengths -------------------
s = prs.slides.add_slide(blank)
set_bg(s, IVORY)
add_header_bar(s, "LinkedIn audit  ·  Part 1 of 2  ·  What it gets right")
add_headline(s, "LinkedIn built the best learning and interview infrastructure in the industry.")

add_bullets(
    s, Inches(0.5), Inches(2.5), Inches(12.3), Inches(4.0),
    [
        {"lead": "20,000+ Learning videos",
         "rest": "the AI coach can lean on a content library no new entrant can match."},
        {"lead": "300+ interview questions, behavioral and technical.",
         "rest": "Interview Prep gives immediate practice and a feedback loop."},
        {"lead": "Quantitative feedback (pace, filler words, buzzwords)",
         "rest": "is actionable in a way peer feedback rarely is."},
        {"lead": "Hundreds of millions of professionals",
         "rest": "the network and content base are a moat."},
    ],
    size=17,
)

# ------------------- SLIDE 5: LinkedIn — weaknesses -------------------
s = prs.slides.add_slide(blank)
set_bg(s, IVORY)
add_header_bar(s, "LinkedIn audit  ·  Part 2 of 2  ·  Where it breaks down")
add_headline(s, "It's a coach for people who already know what they want.")

add_bullets(
    s, Inches(0.5), Inches(2.5), Inches(12.3), Inches(4.2),
    [
        {"lead": "The AI advisor doesn't know your goals or review criteria",
         "rest": "so it can't actually coach you (CareerClimb)."},
        {"lead": "Interview feedback is metric-driven but \"imperfect\"",
         "rest": "it grades pace and fillers, not the subtleties of a real interview (Reddit)."},
        {"lead": "No accountability layer.",
         "rest": "Users still cite \"lack of time\" and want guided next steps beyond a passive course list."},
        {"lead": "AI coaching is buried.",
         "rest": "Nothing kicks in until you enter Learning or Interview Prep — and most of it sits behind a paywall."},
    ],
    size=17,
)

# ------------------- SLIDE 6: Handshake — strengths -------------------
s = prs.slides.add_slide(blank)
set_bg(s, IVORY)
add_header_bar(s, "Handshake audit  ·  Part 1 of 2  ·  What it gets right")
add_headline(s, "Handshake's data integration makes Day 1 surprisingly accurate.")

add_bullets(
    s, Inches(0.5), Inches(2.5), Inches(12.3), Inches(4.2),
    [
        {"lead": "Career inventories",
         "rest": "self-reflection tools students consistently call insightful."},
        {"lead": "Open network with alumni messaging.",
         "rest": "Alumni Hub lowers the barrier to reaching out for help."},
        {"lead": "Tightly integrated onboarding quiz.",
         "rest": "Preferences tune the home feed and alerts immediately."},
        {"lead": "University data piped in (major, grad date)",
         "rest": "job matches are usefully accurate from day one."},
    ],
    size=17,
)

# ------------------- SLIDE 7: Handshake — weaknesses -------------------
s = prs.slides.add_slide(blank)
set_bg(s, IVORY)
add_header_bar(s, "Handshake audit  ·  Part 2 of 2  ·  Where it breaks down")
add_headline(s, "It forces a career label before the student has the self-knowledge to pick one.")

add_bullets(
    s, Inches(0.5), Inches(2.5), Inches(12.3), Inches(4.2),
    [
        {"lead": "The Match Quiz demands a narrow career goal",
         "rest": "students with a general field of interest get pushed into the wrong bucket."},
        {"lead": "Inventories live off-platform.",
         "rest": "CliftonStrengths and FOCUS 2 require external accounts — clunky UX inside a \"unified\" product."},
        {"lead": "Mentorship is underused.",
         "rest": "92% of students say they'd help younger peers, but only ~15% have ever reached out (Handshake)."},
        {"lead": "No active career guidance.",
         "rest": "Handshake itself defers to mentors and career centers for discovery and practical advice."},
    ],
    size=17,
)

# ------------------- SLIDE 8: Synthesis -------------------
s = prs.slides.add_slide(blank)
set_bg(s, NAVY)
add_header_bar(s, "Synthesis", on_dark=True)
add_headline(s, "Here's what's worth saving — and the gap nobody is filling.",
             on_dark=True, size=32)

# Left column: what to keep
add_text(s, Inches(0.5), Inches(2.4), Inches(6), Inches(0.4),
         "KEEP", font=BODY_FONT, size=12, bold=True, color=TEAL)
add_bullets(
    s, Inches(0.5), Inches(2.8), Inches(6.2), Inches(4.0),
    [
        "A large, structured learning library",
        "Job and opportunity recommendation surface",
        "Quantitative feedback on resume and interview reps",
        "User context (school, stage, goals) feeding tailored output",
        "Personalized advising and curricula to develop skills",
    ],
    size=15, color=WHITE, lead_color=TEAL_SOFT, gap_pt=6, line_spacing=1.25,
)

# Right column: the gap
add_rect(s, Inches(7.1), Inches(2.4), Inches(5.75), Inches(4.4), NAVY_DEEP)
add_rect(s, Inches(7.1), Inches(2.4), Inches(0.12), Inches(4.4), TEAL)
add_text(s, Inches(7.4), Inches(2.55), Inches(5.4), Inches(0.4),
         "THE GAP", font=BODY_FONT, size=12, bold=True, color=TEAL)
add_text(s, Inches(7.4), Inches(2.95), Inches(5.4), Inches(3.8),
         "Nobody synthesizes these into a coherent, adaptive experience for a student who doesn't yet know what they want.\n\n"
         "LinkedIn optimizes for the professional who already has a trajectory.\n\n"
         "Handshake optimizes for the employer.\n\n"
         "No one is optimizing for the student at the beginning.",
         font=HEAD_FONT, size=15, color=WHITE, line_spacing=1.35)

# ------------------- SLIDE 9: Introducing ThroughLine (foundation) -------------------
s = prs.slides.add_slide(blank)
set_bg(s, IVORY)
add_header_bar(s, "Introducing ThroughLine  ·  Foundation")
add_headline(s, "Built for the person at the center of the search — and allowed to know who they actually are.")

# Three pillar cards
def pillar(slide, x, title, body):
    add_rect(slide, x, Inches(3.0), Inches(4.0), Inches(3.6), WHITE, RULE)
    add_rect(slide, x, Inches(3.0), Inches(4.0), Inches(0.12), TEAL)
    add_text(slide, x + Inches(0.25), Inches(3.18), Inches(3.6), Inches(0.5),
             title.upper(), font=BODY_FONT, size=11, bold=True, color=TEAL)
    add_text(slide, x + Inches(0.25), Inches(3.55), Inches(3.6), Inches(0.7),
             body["head"], font=HEAD_FONT, size=18, bold=True, color=NAVY,
             line_spacing=1.15)
    add_text(slide, x + Inches(0.25), Inches(4.5), Inches(3.6), Inches(2.0),
             body["text"], font=BODY_FONT, size=13, color=SLATE, line_spacing=1.35)

pillar(s, Inches(0.5), "Pillar 01",
       {"head": "Persistent-memory agent",
        "text": "One model that compresses every interaction into a user-identity dossier. Stores interests, goals, skill gaps, pivots, self-described strengths, past advice. Never loses context across sessions."})
pillar(s, Inches(4.67), "Pillar 02",
       {"head": "Appearance-blind profile",
        "text": "An internal profile, explicitly separate from any external one. Holds what users would never put on a resume: uncertainty, skills in progress, industries they're only curious about."})
pillar(s, Inches(8.83), "Pillar 03",
       {"head": "Factual / advisory disclosure",
        "text": "Every AI output is tagged: sourced fact vs. model judgment. The user always knows when they're getting evidence and when they're getting an opinion."})

# Tagline
add_text(s, Inches(0.5), Inches(2.45), Inches(12.3), Inches(0.5),
         "A persistent-memory agent that delivers tailored advice, curricula, and opportunities — calibrated to a user it genuinely knows.",
         font=HEAD_FONT, size=14, color=MUTED, line_spacing=1.3)

# ------------------- SLIDE 10: Core utilities -------------------
s = prs.slides.add_slide(blank)
set_bg(s, IVORY)
add_header_bar(s, "ThroughLine  ·  Core utilities")
add_headline(s, "From knowing the user to actually moving them forward.")

# two-column feature cards
def feat(slide, x, y, w, h, title, body):
    add_rect(slide, x, y, w, h, WHITE, RULE)
    add_rect(slide, x, y, Inches(0.12), h, TEAL)
    add_text(slide, x + Inches(0.25), y + Inches(0.18), w - Inches(0.5), Inches(0.5),
             title, font=HEAD_FONT, size=18, bold=True, color=NAVY)
    add_text(slide, x + Inches(0.25), y + Inches(0.75), w - Inches(0.5), h - Inches(0.9),
             body, font=BODY_FONT, size=13, color=SLATE, line_spacing=1.35)

feat(s, Inches(0.5), Inches(2.45), Inches(6.15), Inches(4.3),
     "Skill Gap Tracker",
     "User declares a target role, industry, or area of interest. The agent maps current demonstrated skills against that target and tracks the gap closing over time as the user logs new experiences, courses, and projects.\n\nUpdates automatically as the user reports new activity — the bar moves without being asked.")
feat(s, Inches(6.7), Inches(2.45), Inches(6.15), Inches(4.3),
     "Structured Action Planning",
     "Given a high-level goal, the agent generates a phased, time-aware plan with concrete weekly deliverables — ordered by dependency and grounded in the user's actual constraints (hours, baseline, budget, competing commitments).\n\nThe plan is living. When it goes off-track, the agent proposes specific adjustments instead of just flagging the deviation.")

# ------------------- SLIDE 11: AI-native bets -------------------
s = prs.slides.add_slide(blank)
set_bg(s, IVORY)
add_header_bar(s, "ThroughLine  ·  AI-native bets")
add_headline(s, "Two features only a context-aware agent can pull off.")

feat(s, Inches(0.5), Inches(2.45), Inches(6.15), Inches(4.3),
     "Opportunity Discovery Engine",
     "An agentic web research layer that scans against the user's internal context across four categories: paid programs, free resources, student orgs at the user's specific school, and selective unpaid experiences (fellowships, research, competitions).\n\nIt reasons before surfacing — classifies signal vs. noise, explains why, and flags prerequisites the user hasn't met yet.")
feat(s, Inches(6.7), Inches(2.45), Inches(6.15), Inches(4.3),
     "Skill Translator",
     "Reads the user's actual outputs — GitHub repos, research abstracts, project descriptions, written work — and reframes demonstrated competencies into the vocabulary of a target industry.\n\nA philosophy major applying to consulting doesn't know how to reframe their thesis as structured problem-solving. The translator does it automatically, with reasoning shown.")

# ------------------- SLIDE 12: Initial user experience -------------------
s = prs.slides.add_slide(blank)
set_bg(s, NAVY)
add_header_bar(s, "User flow  ·  First five minutes", on_dark=True)
add_headline(s, "Onboarding earns trust by making the user feel seen before it asks for anything.",
             on_dark=True, size=28)

# Step strip
step_y = Inches(3.4)
step_h = Inches(2.6)
gap = Inches(0.2)
step_w = Inches((13.333 - 1.0 - 0.6) / 4)  # 4 steps, padding 0.5 each side, 3 gaps

def step(slide, idx, x, title, text):
    add_rect(slide, x, step_y, step_w, step_h, NAVY_DEEP, NAVY_DEEP)
    add_rect(slide, x, step_y, step_w, Inches(0.45), TEAL)
    add_text(slide, x + Inches(0.2), step_y + Inches(0.05), step_w - Inches(0.4), Inches(0.4),
             f"STEP 0{idx}", font=BODY_FONT, size=11, bold=True, color=NAVY)
    add_text(slide, x + Inches(0.2), step_y + Inches(0.6), step_w - Inches(0.4), Inches(0.6),
             title, font=HEAD_FONT, size=15, bold=True, color=WHITE, line_spacing=1.15)
    add_text(slide, x + Inches(0.2), step_y + Inches(1.25), step_w - Inches(0.4), step_h - Inches(1.4),
             text, font=BODY_FONT, size=11, color=TEAL_SOFT, line_spacing=1.3)

x0 = Inches(0.5)
step(s, 1, x0,
     "Appearance-blind prompt",
     "\"Forget your resume — describe your interests, hobbies, and personality in three sentences.\" Plus three self-description words.")
step(s, 2, x0 + step_w + gap,
     "Memory seed",
     "Input is passed to the persistent-memory agent. This becomes the seed the model builds on — explicitly, so the user understands what's being stored.")
step(s, 3, x0 + 2 * (step_w + gap),
     "Career-fit return",
     "Agent immediately surfaces personality traits, where interests intersect with traits, and which sectors are likely fits.")
step(s, 4, x0 + 3 * (step_w + gap),
     "Deepen the model",
     "Targeted follow-ups extend the seed. People love talking about themselves; the flow uses that to drive early engagement.")

# subline
add_text(s, Inches(0.5), Inches(2.45), Inches(12.3), Inches(0.6),
         "Default career platforms ask for a resume first. ThroughLine asks who you are — and then proves the model is listening.",
         font=HEAD_FONT, size=15, color=TEAL_SOFT, line_spacing=1.3)

# ------------------- SLIDE 13: Prioritization 2x2 -------------------
s = prs.slides.add_slide(blank)
set_bg(s, IVORY)
add_header_bar(s, "Prioritization  ·  Impact vs. effort")
add_headline(s, "Where we'd ship, where we'd bet, and what we'd hold.")

# --- Quadrant chart geometry ---
# Plot area
plot_x = Inches(2.2)
plot_y = Inches(2.4)
plot_w = Inches(8.2)
plot_h = Inches(4.4)
plot_x_emu = plot_x; plot_y_emu = plot_y

# Background plot
add_rect(s, plot_x, plot_y, plot_w, plot_h, WHITE, RULE)

# Quadrant labels (faint background fills)
qw = plot_w / 2
qh = plot_h / 2
# Top-left (low effort, high impact) - Ship First (teal soft tint)
add_rect(s, plot_x, plot_y, qw, qh, RGBColor(0xE6, 0xF7, 0xF5))
# Top-right - Big Bets (light navy tint)
add_rect(s, plot_x + qw, plot_y, qw, qh, RGBColor(0xE8, 0xEC, 0xF4))
# Bottom-left - Revisit
add_rect(s, plot_x, plot_y + qh, qw, qh, RGBColor(0xF1, 0xF3, 0xF8))
# Bottom-right - Deprioritize
add_rect(s, plot_x + qw, plot_y + qh, qw, qh, RGBColor(0xFA, 0xEC, 0xE0))

# crosshair
mid_x = plot_x + qw
mid_y = plot_y + qh
add_rect(s, mid_x - Emu(4500), plot_y, Emu(9000), plot_h, MUTED)
add_rect(s, plot_x, mid_y - Emu(4500), plot_w, Emu(9000), MUTED)

# Axes
# Y axis (impact) - left of plot
ax_y = s.shapes.add_connector(1, plot_x, plot_y, plot_x, plot_y + plot_h)
ax_y.line.color.rgb = NAVY; ax_y.line.width = Pt(1.5)
# X axis (effort) - bottom of plot
ax_x = s.shapes.add_connector(1, plot_x, plot_y + plot_h, plot_x + plot_w, plot_y + plot_h)
ax_x.line.color.rgb = NAVY; ax_x.line.width = Pt(1.5)

# Axis labels
add_text(s, plot_x - Inches(1.4), plot_y - Inches(0.05), Inches(1.3), Inches(0.4),
         "IMPACT", font=BODY_FONT, size=11, bold=True, color=NAVY, align=PP_ALIGN.RIGHT)
add_text(s, plot_x - Inches(1.4), plot_y + plot_h - Inches(0.35), Inches(1.3), Inches(0.4),
         "Low", font=BODY_FONT, size=10, color=MUTED, align=PP_ALIGN.RIGHT)
add_text(s, plot_x - Inches(1.4), plot_y - Inches(0.32), Inches(1.3), Inches(0.4),
         "High", font=BODY_FONT, size=10, color=MUTED, align=PP_ALIGN.RIGHT)

add_text(s, plot_x, plot_y + plot_h + Inches(0.05), plot_w, Inches(0.35),
         "EFFORT", font=BODY_FONT, size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
add_text(s, plot_x - Inches(0.05), plot_y + plot_h + Inches(0.05), Inches(1.2), Inches(0.35),
         "Low", font=BODY_FONT, size=10, color=MUTED, align=PP_ALIGN.LEFT)
add_text(s, plot_x + plot_w - Inches(1.15), plot_y + plot_h + Inches(0.05), Inches(1.2), Inches(0.35),
         "High", font=BODY_FONT, size=10, color=MUTED, align=PP_ALIGN.RIGHT)

# Quadrant titles
add_text(s, plot_x + Inches(0.15), plot_y + Inches(0.1), qw - Inches(0.3), Inches(0.4),
         "SHIP FIRST", font=BODY_FONT, size=11, bold=True, color=TEAL)
add_text(s, plot_x + qw + Inches(0.15), plot_y + Inches(0.1), qw - Inches(0.3), Inches(0.4),
         "BIG BETS", font=BODY_FONT, size=11, bold=True, color=NAVY, align=PP_ALIGN.RIGHT)
add_text(s, plot_x + Inches(0.15), plot_y + qh + Inches(0.05), qw - Inches(0.3), Inches(0.4),
         "REVISIT", font=BODY_FONT, size=11, bold=True, color=MUTED)
add_text(s, plot_x + qw + Inches(0.15), plot_y + qh + Inches(0.05), qw - Inches(0.3), Inches(0.4),
         "DEPRIORITIZE", font=BODY_FONT, size=11, bold=True, color=AMBER, align=PP_ALIGN.RIGHT)

# Data points: (effort, impact) on 0..8 grid
features = [
    ("Persistent Memory Agent",      6, 8, "P1"),
    ("Appearance-Blind Profile",     2, 8, "P2"),
    ("Cold-Start Onboarding",        1, 7, "P3"),
    ("Skill Gap Tracker",            6, 7, "P4"),
    ("Structured Action Planning",   4, 7, "P5"),
    ("Opportunity Discovery Engine", 7, 6, "P6"),
    ("Agent-Mediated Warm Intros",   6, 6, "P7"),
    ("Skill Translator",             4, 6, "P8"),
    ("Factual / Advisory Disclosure",1, 6, "P9"),
    ("Career Trajectory Modeling",   3, 5, "P10"),
    ("Accountability System",        3, 4, "P11"),
    ("AI Interview Prep",            8, 6, "P12"),
]

# scale 0..8 to plot inches, with small inset
inset = Inches(0.18)
plot_inner_w = plot_w - inset * 2
plot_inner_h = plot_h - inset * 2

def to_xy(effort, impact):
    # effort 0..8 -> left..right
    px = plot_x + inset + plot_inner_w * (effort / 8.0)
    # impact 0..8 -> bottom..top (invert)
    py = plot_y + inset + plot_inner_h * (1 - impact / 8.0)
    return px, py

# Plot dots with labels in a legend on right.
# Dot must be large enough for two-character codes (P10..P12) at the chosen font size.
dot_r = Inches(0.34)
for name, eff, imp, code in features:
    cx, cy = to_xy(eff, imp)
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL,
                             cx - dot_r / 2, cy - dot_r / 2, dot_r, dot_r)
    dot.fill.solid()
    dot.fill.fore_color.rgb = TEAL if (imp >= 4 and eff <= 4) else (NAVY if imp >= 4 else AMBER)
    dot.line.color.rgb = WHITE
    dot.line.width = Pt(1.25)
    dot.shadow.inherit = False

    # Code label sits in a textbox the same size as the dot, centered both axes,
    # so PP_ALIGN.CENTER + MIDDLE anchor = visual center of the circle.
    tb = s.shapes.add_textbox(cx - dot_r / 2, cy - dot_r / 2, dot_r, dot_r)
    tf = tb.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.line_spacing = 1.0
    r = p.add_run()
    r.text = code
    r.font.name = BODY_FONT
    r.font.size = Pt(9)
    r.font.bold = True
    r.font.color.rgb = WHITE

# Legend
leg_x = Inches(10.7)
leg_y = Inches(2.4)
leg_w = Inches(2.55)
leg_h = Inches(4.4)
add_rect(s, leg_x, leg_y, leg_w, leg_h, WHITE, RULE)
add_text(s, leg_x + Inches(0.2), leg_y + Inches(0.15), leg_w - Inches(0.4), Inches(0.35),
         "FEATURES", font=BODY_FONT, size=10, bold=True, color=TEAL)
tb = s.shapes.add_textbox(leg_x + Inches(0.2), leg_y + Inches(0.55), leg_w - Inches(0.4), leg_h - Inches(0.7))
tf = tb.text_frame; tf.word_wrap = True
for i, (name, eff, imp, code) in enumerate(features):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.line_spacing = 1.18
    p.space_after = Pt(2)
    rc = p.add_run(); rc.text = code + "  "
    rc.font.name = BODY_FONT; rc.font.size = Pt(9); rc.font.bold = True; rc.font.color.rgb = TEAL
    rn = p.add_run(); rn.text = name
    rn.font.name = BODY_FONT; rn.font.size = Pt(9.5); rn.font.color.rgb = NAVY

# ------------------- SLIDE 14: Risks -------------------
s = prs.slides.add_slide(blank)
set_bg(s, IVORY)
add_header_bar(s, "Product risks")
add_headline(s, "Three failure modes we'd design against from day one.")

risks = [
    ("Data & privacy",
     "The internal profile holds things users would never put on a resume. A breach or unexpected data use is uniquely damaging because we explicitly asked users to be vulnerable."),
    ("Cold-start quality",
     "The first prompt is too thin to generate a genuinely useful career-fit analysis, the first output is generic, and engagement dies before the model has enough context to be useful."),
    ("Trust & accuracy",
     "The model gives confidently wrong career advice at a critical decision point — and the factual/advisory layer didn't catch it because the model failed to flag its own uncertainty."),
]
col_w = Inches(4.0)
col_gap = Inches(0.165)
y0 = Inches(2.6)
for i, (t, b) in enumerate(risks):
    x = Inches(0.5) + i * (col_w + col_gap)
    add_rect(s, x, y0, col_w, Inches(4.1), WHITE, RULE)
    add_rect(s, x, y0, col_w, Inches(0.55), AMBER)
    add_text(s, x + Inches(0.25), y0 + Inches(0.1), col_w - Inches(0.4), Inches(0.45),
             f"RISK 0{i+1}", font=BODY_FONT, size=11, bold=True, color=NAVY)
    add_text(s, x + Inches(0.25), y0 + Inches(0.75), col_w - Inches(0.4), Inches(0.6),
             t, font=HEAD_FONT, size=18, bold=True, color=NAVY, line_spacing=1.1)
    add_text(s, x + Inches(0.25), y0 + Inches(1.5), col_w - Inches(0.4), Inches(2.5),
             b, font=BODY_FONT, size=13, color=SLATE, line_spacing=1.4)

# ------------------- SLIDE 15: Success signals -------------------
s = prs.slides.add_slide(blank)
set_bg(s, IVORY)
add_header_bar(s, "Success signals")
add_headline(s, "How we'd know it's working — behaviorally, not vanity metrics.")

signals = [
    ("Memory depth",
     "Average stored context items per active user grows month-over-month — the persistent model is accumulating real knowledge, not stalling at cold-start."),
    ("Skill gap closure rate",
     "% of tracked skill gaps resolved within 90 days — structured action plans are translating into actual progress toward target roles."),
    ("Return frequency",
     "DAU / WAU > 40% — the persistent memory is delivering enough ongoing value to pull users back consistently."),
    ("Engagement depth",
     "Number of stored context items per active user grows month-over-month, indicating compounding value per session rather than flat surface use."),
]
grid_x = Inches(0.5); grid_y = Inches(2.5); cell_w = Inches(6.15); cell_h = Inches(2.05); cgap = Inches(0.18)
for idx, (t, b) in enumerate(signals):
    r_idx, c_idx = divmod(idx, 2)
    x = grid_x + c_idx * (cell_w + cgap)
    y = grid_y + r_idx * (cell_h + cgap)
    add_rect(s, x, y, cell_w, cell_h, WHITE, RULE)
    add_rect(s, x, y, Inches(0.12), cell_h, TEAL)
    add_text(s, x + Inches(0.25), y + Inches(0.15), cell_w - Inches(0.5), Inches(0.45),
             t, font=HEAD_FONT, size=17, bold=True, color=NAVY)
    add_text(s, x + Inches(0.25), y + Inches(0.7), cell_w - Inches(0.5), cell_h - Inches(0.85),
             b, font=BODY_FONT, size=12.5, color=SLATE, line_spacing=1.35)

# ------------------- SLIDE 16: Close -------------------
s = prs.slides.add_slide(blank)
set_bg(s, NAVY_DEEP)

# left teal bar
add_rect(s, 0, 0, Inches(0.35), SLIDE_H, TEAL)

add_text(s, Inches(0.95), Inches(1.6), Inches(8), Inches(0.4),
         "THE ONE THING TO REMEMBER",
         font=BODY_FONT, size=12, bold=True, color=TEAL)

add_text(s, Inches(0.9), Inches(2.05), Inches(12), Inches(2.6),
         "The network is built.\nThe user isn't.",
         font=HEAD_FONT, size=64, bold=True, color=WHITE, line_spacing=1.05)

add_rect(s, Inches(0.95), Inches(4.55), Inches(1.6), Emu(36000), TEAL)

add_text(s, Inches(0.95), Inches(4.75), Inches(11.5), Inches(2.0),
         "LinkedIn and Handshake approach the search from the employer's side. ThroughLine works from the user's — "
         "building the student into a formidable applicant in whatever field they choose, because it actually knows them: "
         "their strengths, their weaknesses, and what makes them tick.",
         font=HEAD_FONT, size=18, color=WHITE, line_spacing=1.4)

add_text(s, Inches(0.95), Inches(6.6), Inches(8), Inches(0.4),
         "ThroughLine  ·  thank you  ·  Q&A",
         font=BODY_FONT, size=12, bold=True, color=TEAL_SOFT)

# ------------------- Add footers to every slide -------------------
TOTAL = len(prs.slides)
for i, slide in enumerate(prs.slides):
    # Detect dark vs. light: read background fill
    on_dark = False
    try:
        rgb = slide.background.fill.fore_color.rgb
        # Treat NAVY/NAVY_DEEP as dark
        if str(rgb) in (str(NAVY), str(NAVY_DEEP)):
            on_dark = True
    except Exception:
        on_dark = False
    # Skip footer on title (slide 1) and close (slide 16) for visual quiet
    if i == 0 or i == TOTAL - 1:
        continue
    add_footer(slide, i + 1, TOTAL, on_dark=on_dark)

out_path = "/Users/jacobklugerman/Python/ThroughLine_Deck.pptx"
prs.save(out_path)
print(f"Saved {TOTAL} slides → {out_path}")
